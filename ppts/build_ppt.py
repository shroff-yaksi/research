#!/usr/bin/env python3
"""
Build RE-TabSyn_AMTICS_Presentation_v2.pptx
– formatting matches RE-TabSyn_AMTICS_Presentation.pptx (32pt title, 24pt bullets)
– full literature review with all 151 citations by category (Annexure)
– 3-gap problem statement aligned with thesis Chapter 3
– one comprehensive results table (replaces three partial slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1A, 0x29, 0x5E)
LBLUE = RGBColor(0xD6, 0xDC, 0xF0)
MBLUE = RGBColor(0xB8, 0xCC, 0xE8)   # totals / average rows
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY  = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Low-level helpers ─────────────────────────────────────────────────────────

def textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))

def para(tf, text, size, bold=False, italic=False, align=PP_ALIGN.LEFT, color=BLACK,
         space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.alignment    = align
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    r = p.add_run()
    r.text            = text
    r.font.size       = Pt(size)
    r.font.bold       = bold
    r.font.italic     = italic
    r.font.color.rgb  = color
    return p

def first_para(tf, text, size, bold=False, italic=False, align=PP_ALIGN.LEFT, color=BLACK):
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return p

# ── Slide-building helpers ────────────────────────────────────────────────────

def add_title(slide, text, size=32, top=0.30, left=0.50, width=9.0, height=1.25):
    tb = textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first_para(tf, text, size, bold=True, align=PP_ALIGN.CENTER)


def add_subtitle(slide, text, size=24, top=1.75):
    tb = textbox(slide, 0.50, top, 9.0, 0.50)
    tf = tb.text_frame
    tf.word_wrap = True
    first_para(tf, text, size, align=PP_ALIGN.LEFT)


def add_bullets(slide, items, top=1.75, size=24, left=0.50, width=9.0):
    tb = textbox(slide, left, top, width, 7.5 - top - 0.55)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
            p.space_before = Pt(4)
            p.space_after  = Pt(3)
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after  = Pt(3)
        r = p.add_run()
        r.text           = f"•  {item}"
        r.font.size      = Pt(size)
        r.font.color.rgb = BLACK


def add_footnote(slide, text, top, left=0.25, width=9.5):
    tb = textbox(slide, left, top, width, 0.50)
    tf = tb.text_frame
    tf.word_wrap = True
    first_para(tf, text, 9.5, italic=True, color=GREY)


# ── Table builder ─────────────────────────────────────────────────────────────

def build_table(slide, headers, rows, col_widths,
                top=1.75, left=0.25, width=9.5,
                font_size=13, special_last=False):
    """
    special_last=True → last data row gets MBLUE background and bold text
                         (used for Average / Total rows).
    """
    nr = len(rows) + 1       # +1 for header
    nc = len(headers)

    avail_h  = 7.5 - top - 0.55
    row_h    = min(avail_h / nr, 0.52)
    tbl_h    = row_h * nr

    shape = slide.shapes.add_table(
        nr, nc,
        Inches(left), Inches(top),
        Inches(width), Inches(tbl_h))
    tbl = shape.table

    total_w = sum(col_widths)
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = int(Inches(width) * cw / total_w)

    def fmt_cell(cell, text, bg, bold=False, size=None):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r  = p.add_run()
        r.text           = text
        r.font.bold      = bold
        r.font.size      = Pt(size or font_size)
        r.font.color.rgb = WHITE if bg == NAVY else BLACK

    # Header
    for j, h in enumerate(headers):
        fmt_cell(tbl.cell(0, j), h, NAVY, bold=True)

    # Data rows
    for i, row in enumerate(rows):
        is_last = special_last and (i == len(rows) - 1)
        if is_last:
            bg = MBLUE
        elif i % 2 == 0:
            bg = LBLUE
        else:
            bg = WHITE
        for j, val in enumerate(row):
            fmt_cell(tbl.cell(i + 1, j), str(val), bg, bold=is_last)

    return shape, top + tbl_h          # return bottom edge for footnote placement


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = prs.slides.add_slide(BLANK)

def ctr(slide, text, top, height, size, bold=False, italic=False):
    tb = textbox(slide, 0.5, top, 9.0, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first_para(tf, text, size, bold=bold, italic=italic, align=PP_ALIGN.CENTER)

ctr(s, "RE-TabSyn: Controllable Rare-Event Synthetic Data Generation\n"
       "for Financial Tabular Data via Classifier-Free Guidance",
    1.0, 1.4, 22, bold=True)
ctr(s, "Research Presentation",     2.7, 0.5, 16, bold=True, italic=True)
ctr(s, "Submitted by",              3.4, 0.35, 13)
ctr(s, "Yaksi Ketan Shroff\n(202203103510201)", 3.75, 0.75, 14, bold=True)
ctr(s, "under the supervision of:", 4.65, 0.35, 13)
ctr(s, "Dr. Vishvajit Bakrola",     5.0,  0.4,  14, bold=True)
ctr(s, "Department of Computer Science and Engineering, AMTICS, UTU",
    5.4, 0.55, 13)
ctr(s, "April 2026",                6.9,  0.4,  13, bold=True)


# =============================================================================
# SLIDE 2 — Contents
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Contents")
add_bullets(s, [
    "Introduction",
    "Problem Statement",
    "Objectives",
    "Literature Review",
    "Proposed Methodology",
    "Implementation Details",
    "Results & Discussion",
    "Comprehensive Results",
    "Future Work",
    "Conclusion",
    "References",
], size=20)


# =============================================================================
# SLIDE 3 — Introduction
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Introduction")
add_bullets(s, [
    "Synthetic data: artificially generated records that statistically match real data without exposing individuals",
    "Gartner projects synthetic data will overshadow real data in AI training by 2030",
    "Financial sector: privacy-regulated (GDPR, GLBA, PCI-DSS); rare events — fraud, default, bankruptcy — cause billions in losses",
    "Rare events are the most valuable to detect, yet form less than 5% of records",
    "Deep generative models (GANs, VAEs, diffusion) advanced for images and text but lag for mixed-type tabular data",
])


# =============================================================================
# SLIDE 4 — Problem Statement
# Structured around the three research gaps identified in Chapter 3 (Thesis)
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Problem Statement")
add_bullets(s, [
    "Financial datasets severely imbalanced — fraud 0.1–0.5%, loan default 5–15%, bankruptcy 3–8%; rare events most valuable yet least represented",
    "Accuracy paradox: ML classifiers score 99% by predicting majority class — minority recall (fraud, default, bankruptcy) near zero; SMOTE creates unrealistic interpolated samples",
    "Gap 1 — No inference-time minority ratio control: systematic review of 151 papers across 10 categories finds zero instances of retraining-free class-ratio control; all methods mirror or degrade training distribution",
    "Gap 2 — CFG not applied to tabular domain: image-domain CFG [2] requires continuous data; mixed-type tabular features (categorical + continuous) block direct transfer — CTGAN/TVAE mode collapse (Polish Bankr. F1 = 0.112), TabDDPM invalid categories (KS = 0.770)",
    "Gap 3 — No unified fidelity–privacy–control framework: TabSyn leads fidelity (KS = 0.109, zero controllability); DP-CTGAN optimises privacy; SMOTE optimises balance — no single prior method addresses all three",
    "RE-TabSyn fills all three gaps: VAE (mixed-type encoding) + DiT diffusion + CFG → inference-time ratio control via guidance scale w; KS = 0.163, Minority F1 = 0.472 surpasses real baseline 0.458, DCR > 1.0",
], size=20)


# =============================================================================
# SLIDE 5 — Objectives
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Objectives")
add_bullets(s, [
    "Apply Classifier-Free Guidance (CFG) to tabular synthesis for the first time",
    "Enable inference-time control over synthetic class ratios through a single guidance scale w (no retraining)",
    "Preserve statistical fidelity and empirical privacy comparable to the state-of-the-art TabSyn",
    "Improve minority-class detection (minority F1) beyond the real-data baseline",
    "Evaluate on six financial benchmarks (4.8%–44.5% minority) against four baselines — CTGAN, TVAE, TabDDPM, TabSyn",
])


# =============================================================================
# SLIDE 6 — Literature Review: Capability Comparison
# Positions matched to original PPTX (T=2.60 for table, L=0.40, W=9.20)
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Literature Review — Capability Comparison")
add_subtitle(s, "•  Systematic review: 151 papers across 10 categories")

build_table(s,
    headers=["Capability", "CTGAN", "TVAE", "TabDDPM", "TabSyn", "RE-TabSyn"],
    rows=[
        ["Mixed-type support",         "Yes",  "Yes",     "Partial",  "Yes",    "Yes"],
        ["High fidelity (KS < 0.20)",  "Yes",  "Yes",     "No",       "Yes",    "Yes"],
        ["Stable training",            "No",   "Yes",     "Yes",      "Yes",    "Yes"],
        ["Mode coverage",              "No",   "Partial", "Yes",      "Yes",    "Yes"],
        ["Class control",              "No",   "No",      "No",       "No",     "Yes"],
        ["Post-hoc ratio adjustment",  "No",   "No",      "No",       "No",     "Yes"],
        ["Privacy preservation",       "Weak", "Weak",    "Moderate", "Strong", "Strong"],
    ],
    col_widths=[3.4, 1.1, 1.1, 1.3, 1.2, 1.5],
    top=2.60, left=0.40, width=9.20, font_size=16)


# =============================================================================
# SLIDE 6b-1 — Literature Review: All 151 Citations by Category  (1 / 2)
# Categories A–E from Thesis Annexure (21+21+19+8+21 = 90 papers)
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Literature Review — All 151 Citations by Category  (1 / 2)", size=28)

_, btm = build_table(s,
    headers=["Category", "n", "All Citation Numbers", "Key Representative Works"],
    rows=[
        ["A. Diffusion-Based Models",    "21",
         "[5,7,8,10,13,18–20,32–44]",
         "DDPM [5], CFG [7], TabSyn [8], TabDDPM [10], LDM [13], FinDiff [33], CoDi [35], TabDiff [36]"],
        ["B. Privacy-Focused Models",    "21",
         "[55–62,64–76]",
         "DP-SGD [55], PATE-GAN [56], DP-CTGAN [57], PrivSyn [62], Opacus [73], DP-VAE [68]"],
        ["C. GAN-Based Models",          "19",
         "[3,12,14–16,21–31,63,81,96]",
         "GAN [3], CTGAN [21], WGAN-GP [15], CTAB-GAN+ [31], TabFairGAN [23], DP-CTGAN [57]"],
        ["D. Transformer / LLM Models",   "8",
         "[45,47–51,159,160]",
         "TabNet [49], SAINT [50], DiT [160], Attention [159], REaLTabFormer [47], TabLLM [48]"],
        ["E. Evaluation & Benchmarking", "21",
         "[1,98,100,101,118–126,128–133,157,158]",
         "SDV [98], KS Test [132], XGBoost [157], t-SNE [158], SynthEval [130], SynthCity [129]"],
    ],
    col_widths=[2.55, 0.45, 2.85, 3.65],
    top=1.75, left=0.25, width=9.5, font_size=12)

add_footnote(s,
    "Source: Thesis Annexure — Literature Review Scope.  "
    "Citation numbers correspond to thesis reference list [1]–[168].  "
    "Continued on next slide →",
    top=btm + 0.08)


# =============================================================================
# SLIDE 6b-2 — Literature Review: All 151 Citations by Category  (2 / 2)
# Categories F–J from Thesis Annexure (17+18+16+8+2 = 61 papers)
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Literature Review — All 151 Citations by Category  (2 / 2)", size=28)

_, btm = build_table(s,
    headers=["Category", "n", "All Citation Numbers", "Key Representative Works"],
    rows=[
        ["F. Privacy Attacks & Defence", "17",
         "[77,102–117]",
         "MIA [102,103], Model Inversion [110], Re-ID [112], TAMIS [115], Attribute Inf. [77]"],
        ["G. Domain Applications",       "18",
         "[2,9,11,80,82–94,97]",
         "SMOTE [11], GDPR [2], Rare-Event Survey [9], FinSyn [82], EHR-SAFE [83], RareGraph [87]"],
        ["H. Theoretical Foundations",   "16",
         "[4,6,17,53,134,135,138–144,146,159,161]",
         "VAE [4], Imbalanced Data [6], Copulas [134], Optimal Transport [139], Score-SDE [17]"],
        ["I. Recent Advances",            "8",
         "[36,46,149–152,154,156]",
         "Controllable Diffusion [36], Foundation Models [150,151], Goggle [152], Neural ODE [46]"],
        ["J. Uncategorized",              "2",
         "[95,127]",
         "FedAvg [95], CLIP [127] — misc. supporting citations"],
        ["TOTAL",                        "151",
         "[1]–[168]  (excl. datasets [163]–[168] & this work [99])",
         "RE-TabSyn [99] — Shroff & Bakrola, I2IT 2026 (conference paper, this thesis)"],
    ],
    col_widths=[2.55, 0.45, 2.85, 3.65],
    top=1.75, left=0.25, width=9.5, font_size=12, special_last=True)

add_footnote(s,
    "Note: References [163]–[168] are UCI/Kaggle dataset sources; [99] is this thesis/conference paper.  "
    "All 151 peer-reviewed works are categorised above across the 10 categories in the Thesis Annexure.",
    top=btm + 0.08)


# =============================================================================
# SLIDE 7 — Proposed Methodology: Architecture
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Proposed Methodology — Architecture")
add_bullets(s, [
    "RE-TabSyn is a five-module pipeline — removing any module breaks the framework",
    "Preprocessing: median imputation, label encoding, quantile normalisation",
    "VAE Module: encodes mixed-type x into a smooth 64-dimensional latent z; β = 0.1 prevents posterior collapse",
    "Diffusion Module: 4-layer Diffusion Transformer (DiT) denoises in the latent space over T = 1000 steps",
    "CFG Module: blends conditional and unconditional noise predictions at inference",
    "Evaluation Module: KS, TSTR AUC, Minority F1, DCR across three seeds (42, 123, 456)",
])


# =============================================================================
# SLIDE 8 — Proposed Methodology: CFG
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Proposed Methodology — Classifier-Free Guidance")
add_bullets(s, [
    "Training: drop label y with p_uncond = 0.10 and replace it with null token ∅",
    "Model jointly learns conditional ε_cond(z, t, y) and unconditional ε_uncond(z, t, ∅)",
    "Inference blend:  ε̃ = ε_uncond + w · (ε_cond − ε_uncond)",
    "Default guidance scale w = 2.0 drives minority ratio toward ≈ 50%",
    "Changing w at inference adjusts the generated class ratio — no retraining required",
    "First adaptation of CFG (proven in image synthesis) to structured tabular data",
])


# =============================================================================
# SLIDE 9 — Implementation Details: Setup
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Implementation Details — Setup and Hyperparameters")
add_bullets(s, [
    "Hardware: MacBook Air (Apple M3, 16 GB unified RAM) — training completes within a few hours",
    "Software: Python 3.13.5, PyTorch 2.9.1, XGBoost 3.1.2",
    "VAE: encoder [256, 128], decoder [128, 256], latent 64-d, ReLU, β = 0.1",
    "DiT: 4 layers, hidden 256, 4 attention heads, T = 1000 steps, linear β (10⁻⁴ → 0.02)",
    "CFG: label dropout p_uncond = 0.10, guidance scale w = 2.0",
    "Training: Adam, lr = 10⁻³, batch 256, 100 epochs per phase, 80/20 split, 3 seeds",
])


# =============================================================================
# SLIDE 10 — Implementation Details: Datasets
# Positions matched to original PPTX (T=2.60, L=0.40, W=9.20)
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Implementation Details — Datasets and Evaluation")
add_subtitle(s, "•  Six financial benchmarks, minority 4.8%–44.5%")

build_table(s,
    headers=["Dataset", "Task", "Min. %", "Samples"],
    rows=[
        ["Adult Income",      "Income > $50K",             "24.8%", "45,222"],
        ["German Credit",     "Credit risk",               "30.0%",  "1,000"],
        ["Bank Marketing",    "Term deposit subscription", "11.3%", "41,188"],
        ["Credit Approval",   "Approval (653 after NA)",   "44.5%",    "653"],
        ["Lending Club",      "Loan default",              "20.0%", "10,000"],
        ["Polish Bankruptcy", "Corporate bankruptcy",       "4.8%",  "5,000"],
    ],
    col_widths=[2.8, 3.6, 1.5, 1.3],
    top=2.60, left=0.40, width=9.20, font_size=16)


# =============================================================================
# SLIDE 11 — Results: Statistical Fidelity (summary)
# Positions matched to original PPTX (T=3.50, L=0.40, W=9.20)
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Results — Statistical Fidelity")
add_bullets(s, [
    "RE-TabSyn KS = 0.163, comparable to CTGAN (0.153)",
    "TabSyn leads at 0.109; TabDDPM fails entirely (KS 0.770)",
    "+0.054 KS gap is the bounded cost of adding controllability",
], top=1.75, size=20)

build_table(s,
    headers=["Model", "Avg KS ↓", "TSTR AUC ↑", "Minority F1 ↑", "Class Control"],
    rows=[
        ["CTGAN",         "0.153", "0.771", "0.371", "No"],
        ["TabDDPM",       "0.770", "0.518", "—",     "No"],
        ["TabSyn",        "0.109", "0.800", "0.430", "No"],
        ["RE-TabSyn",     "0.163", "0.762", "0.472", "Yes"],
        ["Real baseline", "—",     "0.819", "0.458", "—"],
    ],
    col_widths=[2.4, 2.0, 2.2, 2.2, 2.0],
    top=3.50, left=0.40, width=9.20, font_size=15)


# =============================================================================
# SLIDE 12 — Results: Minority F1 per dataset
# Positions matched to original PPTX (T=3.00, L=0.25, W=9.50)
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Results — Minority F1-Score")
add_bullets(s, [
    "Mean Minority F1 = 0.472, beats real baseline 0.458 (+3.1%)",
    "RE-TabSyn leads every dataset; gains largest at extreme imbalance",
], top=1.75, size=20)

build_table(s,
    headers=["Dataset", "Min %", "Real", "CTGAN", "TabSyn", "RE-TabSyn"],
    rows=[
        ["Adult Income",      "24.8%", "0.543", "0.482", "0.518", "0.552"],
        ["German Credit",     "30.0%", "0.485", "0.425", "0.462", "0.495"],
        ["Bank Marketing",    "11.3%", "0.425", "0.312", "0.385", "0.445"],
        ["Credit Approval",   "44.5%", "0.612", "0.568", "0.595", "0.618"],
        ["Lending Club",      "20.0%", "0.398", "0.325", "0.372", "0.415"],
        ["Polish Bankruptcy",  "4.8%", "0.285", "0.112", "0.245", "0.305"],
        ["Mean",              "—",     "0.458", "0.371", "0.430", "0.472"],
        ["Δ vs. Real",        "—",  "baseline", "−0.087", "−0.028", "+0.014"],
    ],
    col_widths=[2.8, 1.15, 1.35, 1.35, 1.35, 1.6],
    top=3.00, left=0.25, width=9.50, font_size=14, special_last=True)


# =============================================================================
# SLIDE 13 — Results: Class Control, Privacy, Utility
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Results — Class Control, Privacy and Utility")
add_bullets(s, [
    "Minority-ratio control: average achieved 48.1% vs 50% target (1.1% deviation)",
    "Polish Bankruptcy: 4.8% → 50.0% — a tenfold boost with zero seed variance",
    "Privacy (DCR): mean > 1.0 on all datasets — no systematic memorisation",
    "Utility: TSTR AUC 0.762 = 93.0% of the real-data baseline (TabSyn 97.7%)",
    "Guidance scale w is smooth and monotonic: w=0 mirrors training, w=2 balances, w≥3 over-conditions",
], size=20)


# =============================================================================
# SLIDE 13b — Comprehensive Results: All Metrics, All Datasets
# Replaces the three partial result slides (13b, 13c, 13d) from prior revision.
# Matches Table 6.7 (comprehensive_results) from the thesis Research Outcomes chapter.
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Comprehensive Results — All Metrics, All Datasets  (Table 6.7)", size=28)

_, btm = build_table(s,
    headers=["Dataset",       "Min%",
             "CTGAN KS↓",    "TabSyn KS↓", "RE KS↓",
             "Real F1↑",     "CTGAN F1↑",  "TabSyn F1↑", "RE F1↑",
             "Ach%",         "Δ"],
    rows=[
        ["Adult",        "24.8%", "0.152", "0.098", "0.152", "0.543", "0.482", "0.518", "0.552", "49.6%", "0.4%"],
        ["German Cr.",   "30.0%", "0.145", "0.112", "0.156", "0.485", "0.425", "0.462", "0.495", "44.8%", "5.2%"],
        ["Bank Mktg.",   "11.3%", "0.178", "0.115", "0.211", "0.425", "0.312", "0.385", "0.445", "50.2%", "0.2%"],
        ["Credit Appr.", "45.0%", "0.165", "0.125", "0.209", "0.612", "0.568", "0.595", "0.618", "49.6%", "0.4%"],
        ["Lending Club", "20.0%", "0.138", "0.095", "0.140", "0.398", "0.325", "0.372", "0.415", "50.1%", "0.1%"],
        ["Polish Bankr.", "4.8%", "0.142", "0.108", "0.108", "0.285", "0.112", "0.245", "0.305", "50.0%", "0.0%"],
        ["Average",      "22.5%", "0.153", "0.109", "0.163", "0.458", "0.371", "0.430", "0.472", "48.1%", "1.1%"],
    ],
    col_widths=[1.85, 0.60, 0.82, 0.87, 0.77, 0.77, 0.82, 0.87, 0.82, 0.80, 0.57],
    top=1.75, left=0.25, width=9.5, font_size=11, special_last=True)

add_footnote(s,
    "KS↓ lower is better; F1↑ higher is better; Ach% = achieved minority ratio at w=2.0; "
    "Δ = absolute deviation from 50% target.  RE = RE-TabSyn.  "
    "TSTR AUC: RE-TabSyn 0.762 (93.0% of real 0.819); TabSyn 0.800 (97.7%); CTGAN 0.771 (94.1%).  "
    "DCR mean > 1.0 on all datasets — no systematic memorisation.  TabDDPM avg KS = 0.770 (non-viable, omitted per-dataset).",
    top=btm + 0.08)


# =============================================================================
# SLIDE 14 — Future Work
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Future Work")
add_bullets(s, [
    "Near-term: DDIM fast sampling (10–20× speed-up); DP-SGD for formal (ε, δ)-differential privacy (ε = 2.73 feasibility shown)",
    "Medium-term: multi-class CFG for credit ratings (AAA–D); formal membership-inference evaluation; scaled DiT (8–12 layers)",
    "Long-term: federated RE-TabSyn for cross-institution AML; relational multi-table synthesis; fairness-aware multi-attribute conditioning",
    "Domain transfer beyond finance: healthcare records, cybersecurity logs, manufacturing quality control",
    "Production: REST API, constraint-enforcement layer, distribution-drift monitoring, practitioner user study",
])


# =============================================================================
# SLIDE 15 — Conclusion
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Conclusion")
add_bullets(s, [
    "RE-TabSyn is the first tabular synthesis framework with inference-time class-ratio control via CFG",
    "Minority representation boosted from training ratios of 4.8–44.5% to ≈ 50% with 1.1% average deviation",
    "Minority F1 = 0.472 on balanced synthetic data surpasses the real-data baseline of 0.458 (+3.1%) across all six datasets",
    "Empirical privacy preserved: mean DCR > 1.0 on every dataset; no systematic memorisation",
    "Fidelity trade-off is bounded: KS = 0.163 vs TabSyn 0.109 — comparable to CTGAN (0.153) while adding control",
    "Core findings accepted at the I2IT International Conference",
])


# =============================================================================
# SLIDE 16 — References
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "References")
refs = [
    '[1]  J. Ho, A. Jain and P. Abbeel, "Denoising Diffusion Probabilistic Models," NeurIPS, 2020.',
    '[2]  J. Ho and T. Salimans, "Classifier-Free Diffusion Guidance," NeurIPS Workshop, 2021.',
    '[3]  H. Zhang et al., "Mixed-Type Tabular Data Synthesis with Score-based Diffusion (TabSyn)," ICLR, 2024.',
    '[4]  L. Xu et al., "Modeling Tabular Data using Conditional GAN (CTGAN)," NeurIPS, 2019.',
    '[5]  A. Kotelnikov et al., "TabDDPM: Modelling Tabular Data with Diffusion Models," ICML, 2023.',
    '[6]  D. P. Kingma and M. Welling, "Auto-Encoding Variational Bayes," ICLR, 2014.',
    '[7]  R. Rombach et al., "High-Resolution Image Synthesis with Latent Diffusion Models," CVPR, 2022.',
    '[8]  N. V. Chawla et al., "SMOTE: Synthetic Minority Over-sampling Technique," JAIR, 2002.',
    '[9]  M. Abadi et al., "Deep Learning with Differential Privacy," ACM CCS, 2016.',
    '[10] Y. K. Shroff, "Controllable Rare Event Synthetic Data Generation for Financial Tabular Data via CFG," I2IT International Conference, 2026.',
]
tb = textbox(s, 0.50, 1.75, 9.0, 4.95)
tf = tb.text_frame
tf.word_wrap = True
for i, ref in enumerate(refs):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run()
    r.text           = ref
    r.font.size      = Pt(13.5)
    r.font.color.rgb = BLACK


# =============================================================================
# SLIDE 17 — Thank You
# =============================================================================
s = prs.slides.add_slide(BLANK)
tb = textbox(s, 1.0, 2.9, 8.0, 1.2)
tf = tb.text_frame
p  = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r  = p.add_run()
r.text           = "Thank You"
r.font.size      = Pt(40)
r.font.color.rgb = BLACK


# =============================================================================
# Save
# =============================================================================
OUT = "/Users/shroffyaksi/Desktop/Research/ppts/RE-TabSyn_AMTICS_Presentation_v2.pptx"
prs.save(OUT)
print(f"Saved  : {OUT}")
print(f"Slides : {len(prs.slides)}")
for i, sl in enumerate(prs.slides, 1):
    texts = [sh.text_frame.text[:50].strip() for sh in sl.shapes if sh.has_text_frame]
    print(f"  {i:2d}. {texts[0] if texts else '(blank)'}")
